from pptx import Presentation
import sys

def read_pptx(filename):
    try:
        prs = Presentation(filename)
        with open("ppt_content.txt", "w", encoding="utf-8") as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
                f.write("\n")
    except Exception as e:
        print(f"Error reading {filename}: {e}")

if __name__ == "__main__":
    read_pptx("Project_Presentation_Template.pptx")
